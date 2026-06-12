"""
document_processor.py
Downloads a file from a URL and processes it into chunks.
Supports: PDF, DOC/DOCX, PPT/PPTX
"""

import logging
import os
import re
import tempfile
from pathlib import Path
from typing import List

import httpx

try:
    from langchain_core.documents import Document
except ImportError:
    from langchain.schema import Document

from langchain_text_splitters import RecursiveCharacterTextSplitter
from config.settings import CHUNK_SIZE, CHUNK_OVERLAP, MIN_CHUNK_LENGTH

logger = logging.getLogger("mentor-x-ai.document_processor")

TEXT_SPLITTER = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    separators=["\n\n", "\n", ". ", " ", ""],
    add_start_index=True,
)


# ── Loaders ────────────────────────────────────────────────────

def _load_pdf(path: str) -> List[Document]:
    from langchain_community.document_loaders import PyPDFLoader
    return PyPDFLoader(path).load()


def _load_docx(path: str) -> List[Document]:
    try:
        import docx
        doc = docx.Document(path)
        text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return [Document(page_content=text, metadata={"source": path, "page_number": 1})]
    except ImportError:
        # Fallback: UnstructuredWordDocumentLoader
        from langchain_community.document_loaders import UnstructuredWordDocumentLoader
        return UnstructuredWordDocumentLoader(path).load()


def _load_pptx(path: str) -> List[Document]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        pages = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                pages.append(Document(
                    page_content="\n".join(texts),
                    metadata={"source": path, "page_number": i},
                ))
        return pages
    except ImportError:
        from langchain_community.document_loaders import UnstructuredPowerPointLoader
        return UnstructuredPowerPointLoader(path).load()


LOADERS = {
    ".pdf":  _load_pdf,
    ".docx": _load_docx,
    ".doc":  _load_docx,
    ".pptx": _load_pptx,
    ".ppt":  _load_pptx,
}


# ── Text cleaner ───────────────────────────────────────────────

def _clean_text(text: str) -> str:
    text = re.sub(r"\S+@\S+\.\S+", "", text)
    text = re.sub(r"http\S+|www\.\S+", "", text)
    text = re.sub(r"[∗†‡§¶•]", "", text)
    text = re.sub(r"\b\d+\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" \n|\n ", "\n", text)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    return text.replace("“", '"').replace("”", '"').replace("’", "'").strip()


# ── Main processor ─────────────────────────────────────────────

def _normalize_url(file_url: str) -> tuple[str, str | None]:
    """
    Normalize special-case URLs (Google Drive, Dropbox, etc.) to direct-download links.
    Returns (normalized_url, detected_ext_or_None).
    """
    import re

    # Google Drive: /file/d/FILE_ID/view  →  direct download + assume PDF
    gd_match = re.match(r"https://drive\.google\.com/file/d/([^/]+)", file_url)
    if gd_match:
        file_id = gd_match.group(1)
        direct = f"https://drive.google.com/uc?export=download&id={file_id}"
        logger.info("Google Drive URL detected — converting to direct download: %s", direct)
        return direct, ".pdf"   # Google Drive links are almost always PDFs

    # Dropbox: ?dl=0  →  ?dl=1
    if "dropbox.com" in file_url:
        direct = re.sub(r"[?&]dl=0", "", file_url) + ("&dl=1" if "?" in file_url else "?dl=1")
        return direct, None

    return file_url, None


def process_from_url(
    file_url: str,
    file_name: str,
    material_id: int,
    lesson_id: int | None = None,
) -> List[Document]:
    """
    Download a file from file_url, detect type from file_name,
    and return processed + split Document chunks.

    Raises ValueError if file type is unsupported.
    Raises RuntimeError if download fails.
    """
    # ── Normalize URL (Google Drive, Dropbox, etc.) ────────────
    file_url, hint_ext = _normalize_url(file_url)

    # ── Detect file extension ──────────────────────────────────
    ext = Path(file_name).suffix.lower()
    if ext not in LOADERS and hint_ext:
        ext = hint_ext
    if ext not in LOADERS:
        url_path = file_url.split("?")[0]
        ext = Path(url_path).suffix.lower()
    if ext not in LOADERS:
        raise ValueError(f"Unsupported file type: '{ext}'. Supported: {list(LOADERS)}")

    # ── Download to temp file ──────────────────────────────────
    logger.info("Downloading material_id=%d from %s", material_id, file_url[:80])
    try:
        with httpx.Client(timeout=60, follow_redirects=True) as client:
            response = client.get(file_url)
            response.raise_for_status()
            # Google Drive may return a virus-scan warning page for large files
            content_type = response.headers.get("content-type", "")
            if "text/html" in content_type:
                raise RuntimeError(
                    "Server returned HTML instead of a file. "
                    "For Google Drive: make sure the file is shared publicly (Anyone with the link)."
                )
    except httpx.HTTPError as e:
        raise RuntimeError(f"Failed to download file: {e}") from e

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(response.content)
        tmp_path = tmp.name

    try:
        # ── Load ───────────────────────────────────────────────
        loader_fn = LOADERS[ext]
        raw_pages = loader_fn(tmp_path)
        logger.info("Loaded %d pages from %s", len(raw_pages), file_name)

        # ── Clean ──────────────────────────────────────────────
        cleaned = []
        for page in raw_pages:
            text = _clean_text(page.page_content)
            if len(text.split()) < 20:
                continue
            page.page_content = text
            page.metadata.update({
                "source": file_name,
                "material_id": material_id,
                "lesson_id": lesson_id,
            })
            cleaned.append(page)

        # ── Split ──────────────────────────────────────────────
        chunks = TEXT_SPLITTER.split_documents(cleaned)

        # ── Filter tiny chunks ─────────────────────────────────
        chunks = [c for c in chunks if len(c.page_content.strip()) >= MIN_CHUNK_LENGTH]

        # ── Deduplicate ────────────────────────────────────────
        seen, unique = set(), []
        for chunk in chunks:
            key = chunk.page_content.strip()
            if key not in seen:
                seen.add(key)
                unique.append(chunk)

        logger.info(
            "Processed material_id=%d → %d chunks from '%s'",
            material_id, len(unique), file_name,
        )
        return unique

    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
